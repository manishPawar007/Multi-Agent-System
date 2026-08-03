import os
import zipfile
import pandas as pd
from pathlib import Path
from typing import Dict, Any, Tuple
from backend.app.utils.logger import logger

# PyMuPDF (fitz) support
try:
    import fitz
except ImportError:
    fitz = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

try:
    import easyocr
    EASY_OCR_READER = easyocr.Reader(['en'], gpu=False)
except Exception:
    EASY_OCR_READER = None

class DocumentParser:
    @staticmethod
    def parse_file(file_path: str) -> Tuple[str, Dict[str, Any]]:
        path = Path(file_path)
        ext = path.suffix.lower()
        metadata = {
            "filename": path.name,
            "extension": ext,
            "size_bytes": path.stat().st_size if path.exists() else 0
        }

        try:
            if ext == ".pdf":
                return DocumentParser._parse_pdf(file_path), metadata
            elif ext in [".docx", ".doc"]:
                return DocumentParser._parse_docx(file_path), metadata
            elif ext in [".txt", ".md"]:
                return DocumentParser._parse_txt(file_path), metadata
            elif ext == ".csv":
                return DocumentParser._parse_csv(file_path), metadata
            elif ext in [".xlsx", ".xls"]:
                return DocumentParser._parse_excel(file_path), metadata
            elif ext in [".pptx", ".ppt"]:
                return DocumentParser._parse_pptx(file_path), metadata
            elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".webp"]:
                return DocumentParser._parse_image(file_path), metadata
            elif ext == ".zip":
                return DocumentParser._parse_zip(file_path), metadata
            elif ext in [".mp3", ".wav", ".mp4", ".mkv"]:
                return f"[Media File: {path.name} | Type: {ext}]", metadata
            else:
                return DocumentParser._parse_txt(file_path), metadata
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}")
            return f"Error extracting content: {str(e)}", metadata

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        text_content = []
        # Priority 1: PyMuPDF (fitz)
        if fitz:
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    text_content.append(page.get_text())
                if text_content and "".join(text_content).strip():
                    return "\n\n".join(text_content)
            except Exception as e:
                logger.warning(f"PyMuPDF failed on {file_path}: {e}")

        # Priority 2: pypdf fallback
        if pypdf:
            try:
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_content.append(t)
                if text_content and "".join(text_content).strip():
                    return "\n\n".join(text_content)
            except Exception as e:
                logger.warning(f"PyPDF failed on {file_path}: {e}")

        return "No extractable text found in PDF."

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        if docx:
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return "DOCX parser not installed."

    @staticmethod
    def _parse_txt(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _parse_csv(file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()

    @staticmethod
    def _parse_excel(file_path: str) -> str:
        excel_file = pd.ExcelFile(file_path)
        sheets_text = []
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet)
            sheets_text.append(f"--- Sheet: {sheet} ---\n{df.to_string()}")
        return "\n\n".join(sheets_text)

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        if Presentation:
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))
            return "\n\n".join(slides_text)
        return "PPTX parser not available."

    @staticmethod
    def _parse_image(file_path: str) -> str:
        # Priority 1: Tesseract
        if Image and pytesseract:
            try:
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img)
                if text and text.strip():
                    return text
            except Exception as e:
                logger.warning(f"Tesseract OCR failed: {e}")

        # Priority 2: EasyOCR fallback
        if EASY_OCR_READER:
            try:
                results = EASY_OCR_READER.readtext(file_path, detail=0)
                if results:
                    return "\n".join(results)
            except Exception as e:
                logger.warning(f"EasyOCR failed: {e}")

        return "[Image parsed - no text extracted via OCR engines]"

    @staticmethod
    def _parse_zip(file_path: str) -> str:
        contents = []
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            for file_name in zip_ref.namelist():
                contents.append(f"Archive entry: {file_name}")
        return "\n".join(contents)
